#!/usr/bin/env python3
# 產生「TC-DRMS v5 ADOS 架構對齊工作坊」修訂版 deck
# 內容＝Google Slides 現況版全文 ＋ WORKSHOP_DECK_REVISION.md 修改全套用
# 執行：python3 gen_workshop_deck.py → 輸出 TC-DRMS_v5_ADOS工作坊_修訂版.pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

NAVY = RGBColor(0x16, 0x32, 0x4F)
BLUE = RGBColor(0x1D, 0x4E, 0x89)
SKY = RGBColor(0xB9, 0xCD, 0xE2)
GOLD = RGBColor(0xC0, 0x7F, 0x1F)
RED = RGBColor(0xB6, 0x3A, 0x2E)
GREEN = RGBColor(0x2E, 0x7D, 0x54)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Microsoft JhengHei'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def _run(p, text, size=16, bold=False, color=NAVY, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT)
    return r


def txt(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, runs, align=PP_ALIGN.LEFT, space_after=4, space_before=0, line=1.12):
    p = tf.paragraphs[0] if (len(tf.paragraphs) == 1 and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    for spec in runs:
        _run(p, spec[0], **(spec[1] if len(spec) > 1 else {}))
    return p


def rect(slide, x, y, w, h, fill, line_color=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    return s


def header(slide, kicker, title):
    rect(slide, 0.7, 0.52, 0.5, 0.055, GOLD, shape=MSO_SHAPE.RECTANGLE)
    tf = txt(slide, 0.7, 0.62, 11.9, 1.0)
    para(tf, [(kicker, dict(size=12, bold=True, color=GREY))], space_after=2)
    para(tf, [(title, dict(size=27, bold=True, color=NAVY))])


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add(): return prs.slides.add_slide(BLANK)


# ── 1 封面 ─────────────────────────────────────────────
s = add()
rect(s, 0, 0, W, 7.5, NAVY, shape=MSO_SHAPE.RECTANGLE)
rect(s, 5.97, 2.28, 1.4, 0.06, GOLD, shape=MSO_SHAPE.RECTANGLE)
tf = txt(s, 1.2, 2.55, 10.93, 3.2, MSO_ANCHOR.TOP)
para(tf, [('TC-DRMS v5', dict(size=20, bold=True, color=GOLD))], align=PP_ALIGN.CENTER, space_after=10)
para(tf, [('ADOS 架構對齊工作坊', dict(size=44, bold=True, color=WHITE))], align=PP_ALIGN.CENTER, space_after=12)
para(tf, [('Domain × Event × Decision', dict(size=21, color=SKY))], align=PP_ALIGN.CENTER)
tf = txt(s, 1.2, 6.5, 10.93, 0.5)
para(tf, [('2026-07-09', dict(size=13, color=SKY))], align=PP_ALIGN.CENTER)
notes(s, '開場30秒：各位學長，今天不是來看產品的。我先講一個現場：紙本報到、十幾個 LINE 群、150 多條家訪動線靠人腦記——斷的不是網路，是資訊鏈。我的設計理念只有兩句：第一，平時就在用的系統，災時才用得起來；第二，志工的手機裡不需要多一個 APP。接下來一小時，請各位幫我把這兩句話背後的 Domain、Event、Decision 對齊。')

# ── 2 今天要做什麼 ─────────────────────────────────────
s = add()
header(s, 'AGENDA', '今天要做什麼？')
c = rect(s, 0.7, 1.9, 5.85, 4.6, LIGHT)
tf = c.text_frame
para(tf, [('✕ 今天不討論', dict(size=20, bold=True, color=RED))], space_after=14)
for t in ['產品功能細節', 'UI 畫面與操作流程', '個別系統偏好']:
    para(tf, [('—  ' + t, dict(size=17, color=GREY))], space_after=10)
c = rect(s, 6.78, 1.9, 5.85, 4.6, NAVY)
tf = c.text_frame
para(tf, [('✓ 專注對齊', dict(size=20, bold=True, color=GOLD))], space_after=14)
for en, zh in [('Domain', '領域模型'), ('Event', '關鍵事件'), ('Decision', '智慧決策')]:
    para(tf, [(en + '  ', dict(size=17, bold=True, color=WHITE)), (zh, dict(size=17, color=SKY))], space_after=10)

# ── 3 問題（新增）───────────────────────────────────────
s = add()
header(s, 'WHY NOW', '問題：斷的不是網路，是資訊鏈')
pains = [
    ('紙本報到簿', '人數永遠對不上'),
    ('十幾個 LINE 群', '重要通報被訊息洪水洗掉'),
    ('150+ 條家訪動線', '靠人腦記、靠口頭交接'),
    ('收班交接', '沒有快照，下一梯從零開始'),
]
for i, (t, d) in enumerate(pains):
    x = 0.7 + (i % 2) * 6.08
    y = 1.85 + (i // 2) * 1.62
    c = rect(s, x, y, 5.85, 1.42, LIGHT)
    rect(s, x, y + 0.18, 0.07, 1.06, GOLD, shape=MSO_SHAPE.RECTANGLE)
    tf = c.text_frame
    tf.margin_left = Inches(0.28)
    para(tf, [(t, dict(size=18, bold=True, color=NAVY))], space_after=3)
    para(tf, [(d, dict(size=14, color=GREY))])
c = rect(s, 0.7, 5.42, 11.93, 1.12, NAVY)
tf = c.text_frame
para(tf, [('代價：重複作業、漏接個案、無法回溯', dict(size=18, bold=True, color=WHITE))], align=PP_ALIGN.CENTER)
notes(s, '講者：今天所有架構討論，都是為了修這條鏈。｜「150+ 條家訪動線」出自 LOA_ROLES_SPEC 引台南強震安心家訪紀錄，上台前口頭確認數字。')

# ── 4 我們正在建造什麼 ──────────────────────────────────
s = add()
header(s, 'VISION', '我們正在建造什麼？')
tf = txt(s, 0.7, 1.8, 11.93, 1.4)
para(tf, [('從紀錄管理，走向驅動運作', dict(size=15, color=GREY))], align=PP_ALIGN.CENTER, space_after=8)
para(tf, [('DRMS 紀錄系統', dict(size=27, bold=True, color=GREY)),
          ('  ➜  ', dict(size=27, bold=True, color=GOLD)),
          ('ADOS 作業系統', dict(size=27, bold=True, color=NAVY))], align=PP_ALIGN.CENTER)
rows = [
    ('消除斷層', '打通現場與戰情中心資訊鏈'),
    ('即時調度', '從依靠經驗，轉型為數據驅動'),
    ('事件核心', '以事件回饋自動修正決策'),
]
for i, (t, d) in enumerate(rows):
    c = rect(s, 0.7, 3.5 + i * 1.06, 11.93, 0.88, LIGHT)
    tf = c.text_frame
    para(tf, [('➔ ', dict(size=16, bold=True, color=GOLD)),
              (t + '：', dict(size=16, bold=True, color=NAVY)),
              (d, dict(size=16, color=GREY))])

# ── 5 理念一 平戰一體（新增）────────────────────────────
s = add()
header(s, '設計理念 一', '平戰一體')
tf = txt(s, 0.7, 1.72, 11.93, 0.8)
para(tf, [('「平時沒人用的系統，災時一定沒人會用。」', dict(size=22, bold=True, color=BLUE))], align=PP_ALIGN.CENTER)
cards = [
    ('平時', GREEN, '同一套系統跑日常', '報到、演練、物資盤點'),
    ('災時', RED, '切出獨立事件庫', '任務／個案／金援自包含，不汙染日常資料'),
    ('結案', BLUE, '事件封存 ➜ 演練劇本', '匯出成劇本，回到平時訓練'),
]
for i, (t, col, d1, d2) in enumerate(cards):
    x = 0.7 + i * 4.11
    c = rect(s, x, 2.72, 3.9, 2.5, LIGHT)
    rect(s, x, 2.72, 3.9, 0.14, col, shape=MSO_SHAPE.RECTANGLE)
    tf = c.text_frame
    para(tf, [(t, dict(size=20, bold=True, color=col))], align=PP_ALIGN.CENTER, space_after=8)
    para(tf, [(d1, dict(size=16, bold=True, color=NAVY))], align=PP_ALIGN.CENTER, space_after=4)
    para(tf, [(d2, dict(size=13.5, color=GREY))], align=PP_ALIGN.CENTER)
c = rect(s, 0.7, 5.55, 11.93, 1.0, NAVY)
tf = c.text_frame
para(tf, [('每一次真實救災，都變成下一次的演練教材', dict(size=17, bold=True, color=WHITE)),
          ('（story ➜ drill 閉環）', dict(size=17, color=GOLD))], align=PP_ALIGN.CENTER)
notes(s, '對 IT 講：事件隔離＝多災並行不互汙、稽核乾淨、影響範圍可控。｜被問「切開後怎麼合回來」：master 資料唯讀引用、事件庫只寫事件資料、結案回寫摘要。')

# ── 6 理念二 LINE OA（新增）─────────────────────────────
s = add()
header(s, '設計理念 二', 'LINE OA ＝ 資料路由器')
tf = txt(s, 0.7, 1.72, 11.93, 0.8)
para(tf, [('「志工的手機裡，不需要多一個 APP。」', dict(size=22, bold=True, color=BLUE))], align=PP_ALIGN.CENTER)
rows = [
    ('一個入口', 'LINE 官方帳號——LINE 本身就誕生於 311 震災後的通訊需求'),
    ('六個權限面', '志工／班長／司機／香積／訪視／幹部——不映射組織圖，壓到最低教學成本'),
    ('八個串接點', '報到、接單、回報、叫料、簽收、點名 SOS、交接、結案 ➜ 全部寫入唯一資料源'),
]
for i, (t, d) in enumerate(rows):
    c = rect(s, 0.7, 2.72 + i * 0.98, 11.93, 0.84, LIGHT)
    tf = c.text_frame
    tf.margin_left = Inches(0.25)
    para(tf, [(t, dict(size=16, bold=True, color=NAVY)),
              ('　' + d, dict(size=14.5, color=GREY))])
c = rect(s, 0.7, 5.75, 11.93, 0.9, NAVY)
tf = c.text_frame
para(tf, [('科技的目的：把志工從表單裡解放出來，把時間還給膚慰與關懷', dict(size=16.5, bold=True, color=GOLD))], align=PP_ALIGN.CENTER)
notes(s, '被問「LINE 掛了怎麼辦」：入口可降級（紙本 QR 備援），資料層不依賴 LINE——LINE 只是路由器，不是資料庫。')

# ── 7 物件（五體環繞）───────────────────────────────────
s = add()
header(s, 'DOMAIN', '我們的世界有哪些物件？')


def obj(x, y, w, h, ico, name, d, fill=LIGHT, line=None, tcol=NAVY, big=False):
    c = rect(s, x, y, w, h, fill, line_color=line, line_w=1.6)
    tf = c.text_frame
    para(tf, [(ico + ' ' + name, dict(size=17 if big else 14.5, bold=True, color=tcol))], align=PP_ALIGN.CENTER, space_after=2)
    para(tf, [(d, dict(size=11.5, color=GREY if fill == LIGHT else SKY))], align=PP_ALIGN.CENTER)


obj(5.27, 1.62, 2.8, 0.95, '🤖', 'AI 智能引擎', '排序・建議・預測（Phase 4）', line=GREY)
obj(1.5, 2.75, 2.9, 1.05, '👤', 'Person 人', '報到・接任務・回報')
obj(8.93, 2.75, 2.9, 1.05, '🪖', 'Squad 班組', '班長帶隊・統一接單')
obj(5.07, 3.28, 3.2, 1.25, '📋', 'Task 任務', '五體交匯・最小執行單位', fill=NAVY, tcol=WHITE, big=True)
obj(1.5, 4.15, 2.9, 1.05, '🏠', 'Case 個案', '受理・分診・追蹤')
obj(8.93, 4.15, 2.9, 1.05, '🚛', 'Vehicle 車輛', '派車・掃碼・扣庫存')
obj(3.55, 5.45, 2.9, 0.95, '💰', 'ReliefFund 金援', '五步驟發放・不可刪帳本')
obj(6.88, 5.45, 2.9, 0.95, '🤝', 'Handover 交接', '收班快照・電子簽名')
tf = txt(s, 0.7, 6.6, 11.93, 0.5)
para(tf, [('輸出層：HQ 儀表板・班長 APP・志工 APP（LINE OA）・財務審計・保險・政府報告', dict(size=12.5, color=GREY))], align=PP_ALIGN.CENTER)
notes(s, '貼紙討論預埋這八個物件——從空白開始但不從零開始（出自 ARCH_V2 五鏈）。')

# ── 8 EDA ──────────────────────────────────────────────
s = add()
header(s, 'EVENT', 'Event-Driven Architecture')
tf = txt(s, 0.7, 1.75, 11.93, 0.6)
para(tf, [('ADOS 不是 Workflow，而是', dict(size=16, color=GREY)),
          ('事件觸發、決策回應', dict(size=16, bold=True, color=NAVY)),
          ('的動態系統。', dict(size=16, color=GREY))], align=PP_ALIGN.CENTER)
evs = [('1', 'Disaster Occurred', '災害發生'), ('2', 'Case Reported', '案件通報'),
       ('3', 'Priority Calculated', '優先級計算'), ('4', 'Squad Dispatched', '小隊派遣'),
       ('5', 'Task Completed', '任務完成')]
for i, (n, en, zh) in enumerate(evs):
    x = 0.7 + i * 2.44
    rect(s, x, 2.9, 2.28, 1.7, NAVY if i < 4 else GREEN, shape=MSO_SHAPE.CHEVRON)
    # chevron 內文字用獨立窄框強制換行，字才不會溢出箭身（白字上白底＝隱形）
    tf = txt(s, x + 0.52, 3.02, 1.5, 1.46, MSO_ANCHOR.MIDDLE)
    para(tf, [(n, dict(size=15, bold=True, color=GOLD))], align=PP_ALIGN.CENTER, space_after=2)
    for word in en.split(' '):  # EN 一律拆兩行，五個箭頭視覺一致且不溢出箭身
        para(tf, [(word, dict(size=11, bold=True, color=WHITE))], align=PP_ALIGN.CENTER, space_after=0, line=1.05)
    para(tf, [(zh, dict(size=10.5, color=SKY))], align=PP_ALIGN.CENTER, space_before=2)
c = rect(s, 0.7, 5.1, 11.93, 0.95, LIGHT, line_color=GOLD, line_w=1.2)
tf = c.text_frame
para(tf, [('↺ 事件回饋（閉環）：', dict(size=15.5, bold=True, color=GOLD)),
          ('Task Completed ➜ 重算 Priority、沉澱成演練劇本', dict(size=15.5, color=NAVY))], align=PP_ALIGN.CENTER)

# ── 9 Decision Fabric ──────────────────────────────────
s = add()
header(s, 'DECISION', 'Decision Fabric：系統的核心價值')
cards = [
    ('優先級決策', '評估嚴重性與時效性，自動排定處理順序（A／B／C 級）'),
    ('派遣決策', '根據距離、專業技能與當前負載，推薦最適指派小隊'),
    ('資源決策', '動態計算最優物資與載具配置，確保前線不斷鏈'),
]
for i, (t, d) in enumerate(cards):
    x = 0.7 + i * 4.11
    c = rect(s, x, 2.1, 3.9, 3.4, LIGHT)
    rect(s, x, 2.1, 3.9, 0.14, BLUE, shape=MSO_SHAPE.RECTANGLE)
    tf = c.text_frame
    para(tf, [(t, dict(size=19, bold=True, color=NAVY))], align=PP_ALIGN.CENTER, space_after=10)
    para(tf, [(d, dict(size=14.5, color=GREY))], align=PP_ALIGN.CENTER)

# ── 10 AI 定位 ─────────────────────────────────────────
s = add()
header(s, 'AI', 'AI 的定位：決策夥伴')
c = rect(s, 0.7, 1.9, 5.85, 4.6, NAVY)
tf = c.text_frame
tf.margin_left = Inches(0.25)
para(tf, [('AI 三大設計原則', dict(size=18, bold=True, color=GOLD))], space_after=12)
for t, d in [('1. 不取代 UI', '專注背景邏輯，不做無謂 Chatbot'),
             ('2. 不取代人類控制權', '保留最終裁決權'),
             ('3. 定位為「決策助理」', '最強大的參謀，不是指揮官')]:
    para(tf, [(t, dict(size=16, bold=True, color=WHITE))], space_after=2)
    para(tf, [(d, dict(size=13.5, color=SKY))], space_after=10)
for i, (t, d) in enumerate([('判斷', '解析語音／通報，自動進行案件分流'),
                            ('建議', '推薦派遣小隊、路線與物資包'),
                            ('摘要', '整理複雜情資，形成指揮官 Dashboard')]):
    c = rect(s, 6.78, 1.9 + i * 1.6, 5.85, 1.4, LIGHT)
    tf = c.text_frame
    tf.margin_left = Inches(0.25)
    para(tf, [('➔ ' + t, dict(size=17, bold=True, color=BLUE))], space_after=3)
    para(tf, [(d, dict(size=14, color=GREY))])

# ── 11 架構：現況與目標（改寫）──────────────────────────
s = add()
header(s, 'ARCHITECTURE', '架構：現況與目標')
c = rect(s, 0.7, 1.85, 11.93, 1.15, LIGHT, line_color=GOLD, line_w=1.5)
tf = c.text_frame
para(tf, [('現況（誠實版）　', dict(size=16, bold=True, color=GOLD)),
          ('前端原型完整｜LINE OA 模擬器完整｜真實 Webhook 0%｜尚無正式後端', dict(size=15, color=NAVY))], align=PP_ALIGN.CENTER)
tf = txt(s, 0.7, 3.2, 11.93, 0.5)
para(tf, [('目標', dict(size=16, bold=True, color=NAVY))])
rows = [
    ('Event Mesh 事件流中樞', '承載即時事件流的核心'),
    ('Digital Twin 數位對照', '人、資源、地理位置的即時關係'),
    ('AI Agent Layer 決策助理層', '判斷、建議、摘要'),
]
for i, (t, d) in enumerate(rows):
    c = rect(s, 0.7, 3.7 + i * 0.92, 11.93, 0.78, NAVY)
    tf = c.text_frame
    tf.margin_left = Inches(0.25)
    para(tf, [(t, dict(size=15.5, bold=True, color=WHITE)),
              ('　—　' + d, dict(size=14, color=SKY))])
tf = txt(s, 0.7, 6.55, 11.93, 0.6)
para(tf, [('今天的對齊，就是從現況走到目標的第一步：API 規格、服務邊界、Schema',
           dict(size=15, bold=True, color=GOLD))], align=PP_ALIGN.CENTER)
notes(s, '現況自己先講＝主導權在你；被問出來就變「原來是空的」。')

# ── 12 演進路線圖 ──────────────────────────────────────
s = add()
header(s, 'ROADMAP', 'ADOS 演進路線圖')
cols = [
    ('已上線', GREEN, '報到・個案・車輛\n任務・HQ 儀表板'),
    ('Phase 0-A', BLUE, 'Squad 班組 schema\n班長行動端 5 顆按鈕'),
    ('Phase 1-D', BLUE, 'Handover 交接快照\n電子簽名・不可修改'),
    ('Phase 3-A', BLUE, 'ReliefFund 五步驟\n不可刪帳本'),
    ('Phase 4', GOLD, 'AI 引擎\n（需 3+ 次真實出班資料）'),
]
for i, (t, col, d) in enumerate(cols):
    x = 0.7 + i * 2.44
    c = rect(s, x, 2.3, 2.28, 3.0, LIGHT)
    rect(s, x, 2.3, 2.28, 0.14, col, shape=MSO_SHAPE.RECTANGLE)
    tf = c.text_frame
    para(tf, [(t, dict(size=16, bold=True, color=col))], align=PP_ALIGN.CENTER, space_after=8)
    for seg in d.split('\n'):
        para(tf, [(seg, dict(size=12.5, color=GREY))], align=PP_ALIGN.CENTER, space_after=3)
tf = txt(s, 0.7, 5.7, 11.93, 0.5)
para(tf, [('每一階段都以「真實出班」驗證後才進下一階', dict(size=14, bold=True, color=NAVY))], align=PP_ALIGN.CENTER)
notes(s, '本頁依 ARCH_V2 HEALTH 現況重建；如你原本的路線圖更完整，直接替換本頁。')

# ── 13 今日核心目標 ────────────────────────────────────
s = add()
header(s, 'GOALS', '今日核心目標')
tf = txt(s, 0.7, 1.7, 11.93, 0.5)
para(tf, [('工作坊結束前，團隊必須共同明確回答這四個核心問題：', dict(size=15, color=GREY))])
qs = ['Domain 定義對嗎？', 'Event 完整嗎？', 'Decision 權責清楚嗎？', '下個 Sprint 做什麼？']
for i, q in enumerate(qs):
    x = 0.7 + (i % 2) * 6.08
    y = 2.4 + (i // 2) * 2.1
    c = rect(s, x, y, 5.85, 1.9, LIGHT)
    tf = c.text_frame
    tf.margin_left = Inches(0.3)
    para(tf, [('0' + str(i + 1), dict(size=26, bold=True, color=GOLD))], space_after=4)
    para(tf, [(q, dict(size=19, bold=True, color=NAVY))])

# ── 14 會後行動 ────────────────────────────────────────
s = add()
header(s, 'NEXT ACTIONS', '會後行動：實踐共享語言')
rows = [
    '撰寫 Architecture Decision Record（ADR）',
    '定義 API 規格與服務邊界',
    '設計符合作業系統願景的 Database Schema',
    '快速產出 Prototype 概念驗證',
]
for i, t in enumerate(rows):
    c = rect(s, 0.7, 2.0 + i * 1.15, 11.93, 0.95, LIGHT)
    tf = c.text_frame
    tf.margin_left = Inches(0.25)
    para(tf, [(str(i + 1) + '　', dict(size=18, bold=True, color=GOLD)),
              (t, dict(size=17, bold=True, color=NAVY))])

# ── 15 ADOS Principle ──────────────────────────────────
s = add()
rect(s, 0, 0, W, 7.5, NAVY, shape=MSO_SHAPE.RECTANGLE)
tf = txt(s, 1.2, 2.6, 10.93, 2.4)
para(tf, [('ADOS Principle', dict(size=16, bold=True, color=GOLD))], align=PP_ALIGN.CENTER, space_after=16)
para(tf, [('「不打造另一個管理系統，', dict(size=30, bold=True, color=WHITE))], align=PP_ALIGN.CENTER, space_after=6)
para(tf, [('而是打造災害應變的作業系統。」', dict(size=30, bold=True, color=WHITE))], align=PP_ALIGN.CENTER)

# ── 16 附錄：戰時授權差分 ───────────────────────────────
s = add()
header(s, '附錄・備而不發', '戰時授權——差分，不是複製')
c = rect(s, 0.7, 1.85, 11.93, 0.95, NAVY)
tf = c.text_frame
para(tf, [('授權矩陣 ＝ 角色 × 動作 × 金額級距 × 情境（平時／戰時）', dict(size=17, bold=True, color=WHITE))], align=PP_ALIGN.CENTER)
tf = txt(s, 0.7, 3.0, 11.93, 0.5)
para(tf, [('戰時不是全面放寬——有放有收：', dict(size=16, bold=True, color=NAVY))])
rows = [
    ('放', GREEN, '組長核准上限提高（幅度今天表決）＋全案標記「戰時放寬」＋48h 補審'),
    ('收', RED, 'master 資料鎖唯讀（單一資訊流）'),
    ('收', RED, '解除戰時比進入更嚴（僅 admin）——避免誤觸清空戰時狀態'),
]
for i, (t, col, d) in enumerate(rows):
    c = rect(s, 0.7, 3.55 + i * 1.0, 11.93, 0.85, LIGHT)
    tf = c.text_frame
    tf.margin_left = Inches(0.25)
    para(tf, [(t + '　', dict(size=16, bold=True, color=col)),
              (d, dict(size=14.5, color=NAVY))])
tf = txt(s, 0.7, 6.7, 11.93, 0.5)
para(tf, [('＝今日核心目標第 3 題「Decision 權責清楚嗎」要各位表決的東西', dict(size=13.5, color=GREY))], align=PP_ALIGN.CENTER)
notes(s, '備而不發：只在被問「戰時誰能決定什麼」「權限會不會亂」時翻出來。內容出自 AUTH_MATRIX_SPEC §4.5。')

prs.core_properties.title = 'TC-DRMS v5 ADOS 架構對齊工作坊（修訂版）'
prs.core_properties.author = 'Mason'
OUT = 'TC-DRMS_v5_ADOS工作坊_修訂版.pptx'
prs.save(OUT)
print('saved', OUT, '| slides =', len(prs.slides.__iter__.__self__._sldIdLst))
